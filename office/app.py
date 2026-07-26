import io

from fastapi import FastAPI, HTTPException, Request, Response
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from docx import Document
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import cm
from reportlab.pdfgen import canvas
import pdfplumber

app = FastAPI()

XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"


@app.get("/health")
async def health():
    return {"ok": True}


@app.post("/excel/create")
async def create_excel(request: Request):
    payload = await request.json()
    sheet_name = str(payload.get("sheetName") or "Foglio1")[:31]
    rows = payload.get("rows") or []
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = sheet_name
    for row in rows:
        sheet.append(row)
    buffer = io.BytesIO()
    workbook.save(buffer)
    buffer.seek(0)
    return Response(content=buffer.read(), media_type=XLSX_MIME)


@app.post("/excel/read")
async def read_excel(request: Request):
    data = await request.body()
    if not data:
        raise HTTPException(400, "File Excel vuoto.")
    try:
        workbook = load_workbook(io.BytesIO(data), data_only=True)
    except Exception as error:  # noqa: BLE001
        raise HTTPException(400, f"Impossibile leggere il file Excel: {error}") from error
    sheets = []
    for sheet in workbook.worksheets:
        rows = [list(row) for row in sheet.iter_rows(values_only=True, max_row=200)]
        sheets.append({"name": sheet.title, "rows": rows})
    return {"sheets": sheets}


@app.post("/pptx/create")
async def create_pptx(request: Request):
    payload = await request.json()
    slides_data = payload.get("slides") or []
    presentation = Presentation()
    layout = presentation.slide_layouts[1]
    for slide_data in slides_data:
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = str(slide_data.get("title") or "")
        bullets = slide_data.get("bullets") or []
        if bullets:
            body = slide.placeholders[1].text_frame
            body.text = str(bullets[0])
            for bullet in bullets[1:]:
                paragraph = body.add_paragraph()
                paragraph.text = str(bullet)
    buffer = io.BytesIO()
    presentation.save(buffer)
    buffer.seek(0)
    return Response(content=buffer.read(), media_type=PPTX_MIME)


@app.post("/docx/create")
async def create_docx(request: Request):
    payload = await request.json()
    title = payload.get("title")
    sections = payload.get("sections") or []
    document = Document()
    if title:
        document.add_heading(str(title), level=1)
    for section in sections:
        heading = section.get("heading")
        if heading:
            document.add_heading(str(heading), level=2)
        for paragraph in section.get("paragraphs") or []:
            document.add_paragraph(str(paragraph))
    buffer = io.BytesIO()
    document.save(buffer)
    buffer.seek(0)
    return Response(content=buffer.read(), media_type=DOCX_MIME)


def _wrap_text(text, width):
    words = text.split()
    lines = []
    current = ""
    for word in words:
        candidate = f"{current} {word}".strip()
        if len(candidate) <= width:
            current = candidate
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines or [""]


@app.post("/pdf/create")
async def create_pdf(request: Request):
    payload = await request.json()
    title = payload.get("title")
    paragraphs = payload.get("paragraphs") or []
    buffer = io.BytesIO()
    pdf_canvas = canvas.Canvas(buffer, pagesize=A4)
    width, height = A4
    y = height - 2 * cm

    if title:
        pdf_canvas.setFont("Helvetica-Bold", 16)
        pdf_canvas.drawString(2 * cm, y, str(title)[:90])
        y -= 1.2 * cm

    pdf_canvas.setFont("Helvetica", 11)
    for paragraph in paragraphs:
        for line in _wrap_text(str(paragraph), 90):
            if y < 2 * cm:
                pdf_canvas.showPage()
                pdf_canvas.setFont("Helvetica", 11)
                y = height - 2 * cm
            pdf_canvas.drawString(2 * cm, y, line)
            y -= 0.6 * cm
        y -= 0.4 * cm

    pdf_canvas.save()
    buffer.seek(0)
    return Response(content=buffer.read(), media_type="application/pdf")


@app.post("/pdf/extract-text")
async def extract_pdf_text(request: Request):
    data = await request.body()
    if not data:
        raise HTTPException(400, "File PDF vuoto.")
    text_parts = []
    try:
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for page in pdf.pages[:50]:
                text_parts.append(page.extract_text() or "")
    except Exception as error:  # noqa: BLE001
        raise HTTPException(400, f"Impossibile leggere il PDF: {error}") from error
    return {"text": "\n\n".join(text_parts).strip()[:40000]}
